from pptx import Presentation

prs = Presentation()

def add_slide(title, content_lines):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(content_lines):
        if i == 0:
            tf.text = line
        else:
            p = tf.add_paragraph()
            p.text = line

# Slides
add_slide("AI Governance Demo",
          ["Grok vs HHI Governance Workflow",
           "Control before output"])

add_slide("What Grok Does",
          ["Structured responses",
           "Uses disclaimers",
           "Avoids legal liability",
           "Advisory only"])

add_slide("Where Grok Fails",
          ["No Decision Boundaries",
           "No Stop conditions",
           "No enforced escalation",
           "Accountability pushed to user"])

add_slide("What HHI Does",
          ["Defines Decision Boundaries",
           "Enforces Escalation",
           "Applies Stop conditions",
           "Binds Accountability"])

add_slide("Key Difference",
          ["Grok avoids risk",
           "HHI controls risk"])

add_slide("Business Impact",
          ["Prevents bad outputs before they exist",
           "Reduces legal and operational risk",
           "Creates audit-ready systems"])

prs.save("HHI_Governance_Demo.pptx")
